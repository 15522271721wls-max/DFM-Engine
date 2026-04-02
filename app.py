import streamlit as st
from pptx import Presentation
import io
import requests
from zhipuai import ZhipuAI
import docx 
from openai import OpenAI
import json
import streamlit.components.v1 as components
import re

# === 网页基础设置 ===
st.set_page_config(page_title="DFM生成引擎", page_icon="⚙️", layout="centered")

# ==========================================
#           🚨 专属系统防盗门 (拦截器)
# ==========================================
def check_password():
    if st.session_state["password_input"] == "JINYA888":
        st.session_state["authenticated"] = True
    else:
        st.session_state["authenticated"] = False
        st.error("❌ 邀请码错误或已失效，请联系系统管理员获取！")

if "authenticated" not in st.session_state:
    st.session_state["authenticated"] = False

if not st.session_state["authenticated"]:
    st.markdown("<br><br><br>", unsafe_allow_html=True)
    st.markdown("<h1 style='text-align: center; color: #1E3A8A; font-size: 3rem;'>🔒 DFM生成引擎</h1>", unsafe_allow_html=True)
    st.markdown("<p style='text-align: center; color: #6B7280; font-size: 1.2rem; margin-bottom: 2rem;'>内部专属自动化演示系统，请输入授权邀请码解锁</p>", unsafe_allow_html=True)
    
    col1, col2, col3 = st.columns([1, 2, 1])
    with col2:
        st.text_input("邀请码", type="password", key="password_input", label_visibility="collapsed", placeholder="请输入邀请码 (如: 123456)")
        st.button("🔓 验证并解锁系统", on_click=check_password, use_container_width=True)
    st.stop()

# === 高级 CSS 魔法 ===
st.markdown("""
<style>
    html, body, [class*="css"] { font-family: "PingFang SC", "Microsoft YaHei", sans-serif !important; }
    .block-container { max-width: 900px !important; padding-top: 2rem !important; }
    .main-title { text-align: center; font-size: 2.8rem; font-weight: 800; background: linear-gradient(90deg, #1E3A8A, #3B82F6); -webkit-background-clip: text; -webkit-text-fill-color: transparent; margin-bottom: 0.5rem; padding-top: 1rem; }
    .sub-title { text-align: center; color: #6B7280; font-size: 1.1rem; font-weight: 500; margin-bottom: 2rem; }
    .stButton > button, .stDownloadButton > button, div[data-testid="stFormSubmitButton"] > button { background: linear-gradient(90deg, #1E3A8A, #3B82F6) !important; color: white !important; font-size: 1.1rem !important; font-weight: bold !important; border-radius: 8px !important; border: none !important; padding: 0.6rem 0 !important; box-shadow: 0 4px 6px rgba(59, 130, 246, 0.2) !important; transition: all 0.3s ease !important; }
    .stButton > button:hover, .stDownloadButton > button:hover, div[data-testid="stFormSubmitButton"] > button:hover { transform: translateY(-2px); box-shadow: 0 6px 12px rgba(59, 130, 246, 0.4) !important; }
    div[data-testid="InputInstructions"] { display: none !important; }
    div[data-baseweb="select"], div[data-baseweb="select"] input { cursor: pointer !important; }
    section[data-testid="stSidebar"] { min-width: 360px !important; max-width: 400px !important; }
    
    /* 🎯 上传按钮美化 & 垃圾桶绝对隔离 🎯 */
    div[data-testid="stFileUploader"] button[kind="secondary"] { background: linear-gradient(90deg, #1E3A8A, #3B82F6) !important; border: none !important; border-radius: 8px !important; color: transparent !important; position: relative; }
    div[data-testid="stFileUploader"] button[kind="secondary"]::after { content: "浏览本地文件" !important; position: absolute; color: white !important; left: 50%; top: 50%; transform: translate(-50%, -50%); font-weight: bold; white-space: nowrap; }
    div[data-testid="stFileUploader"] button[kind="secondary"]:hover { transform: translateY(-2px); box-shadow: 0 4px 10px rgba(59, 130, 246, 0.3) !important; }
    div[data-testid="stUploadedFile"] button, div[data-testid="stFileUploaderFileData"] ~ div button { background: transparent !important; color: inherit !important; box-shadow: none !important; transform: none !important; }
    div[data-testid="stUploadedFile"] button::after, div[data-testid="stFileUploaderFileData"] ~ div button::after { display: none !important; content: none !important; }
    div[data-testid="stFileUploader"] [data-testid="stMarkdownContainer"] p { font-size: 0 !important; }
    div[data-testid="stFileUploader"] [data-testid="stMarkdownContainer"] p::after { content: "拖拽文件至此区域" !important; font-size: 16px !important; font-weight: 600 !important; color: var(--text-color) !important; }
    div[data-testid="stFileUploader"] small { display: none !important; }
    
    /* 🎯 修改：移除跑马灯特效，改为静态深蓝色 🎯 */
    div[data-testid="stExpander"] summary p { 
        font-weight: 900 !important; 
        font-size: 1.15rem !important; 
        color: #1E3A8A !important; /* 固定为专业深蓝色 */
    }
    
    /* 清除任何背景图片残留 */
    [data-testid="stAppViewContainer"], [data-testid="stSidebar"] { background-image: none !important; background-color: white !important; }
    
    /* 教程指引专属强调色 */
    .tutorial-step { color: #1E3A8A; font-weight: 700; font-size: 1.05rem; }
</style>
""", unsafe_allow_html=True)


# === 后台逻辑函数 ===
def extract_text_content(uploaded_file):
    if not uploaded_file: return ""
    try:
        if uploaded_file.name.lower().endswith(('.md', '.txt')): return uploaded_file.getvalue().decode('utf-8')
        elif uploaded_file.name.lower().endswith('.docx'): return '\n'.join([para.text for para in docx.Document(io.BytesIO(uploaded_file.getvalue())).paragraphs])
    except Exception: return ""
    return ""

def parse_markdown(md_text):
    if not md_text: return []
    
    md_text = re.sub(r'(?<!^)(?=\[(?:封面|目录|过渡|左文右图|左右文字)\])', r'\n---\n', md_text)
    md_text = re.sub(r'。', r'。\n', md_text)
    md_text = re.sub(r'(?=\[AI配图:)', r'\n', md_text)
    
    parsed_slides = []
    for slide in md_text.split('---'):
        slide = slide.strip()
        if not slide: continue
        
        l_type, title, ai_prompt = "左右文字", "", ""
        left_body, right_body = [], []
        current_target = left_body
        
        lines = slide.split('\n')
        for line in lines:
            line_str = line.strip()
            if not line_str: continue
            
            if line_str == '[左]': current_target = left_body; continue
            elif line_str == '[右]': current_target = right_body; continue
            
            match_ai = re.search(r'\[AI配图:(.*?)\]', line_str)
            if match_ai:
                ai_prompt = match_ai.group(1).strip()
                line_str = re.sub(r'\[AI配图:.*?\]', '', line_str).strip()
                if not line_str: continue
                
            match_layout = re.match(r'^\[(封面|目录|过渡|左文右图|左右文字|正文)\](.*)', line_str)
            if match_layout:
                l_type = match_layout.group(1)
                if l_type == "正文": l_type = "左右文字" 
                remainder = match_layout.group(2).strip()
                if remainder:
                    if remainder.startswith('# '): title = remainder[2:]
                    else:
                        split_idx = -1
                        for punct in [' ', '。', '：', ':', '\t']:
                            idx = remainder.find(punct)
                            if idx != -1 and idx < 20: split_idx = idx; break
                        if split_idx != -1:
                            title = remainder[:split_idx].strip()
                            if remainder[split_idx+1:].strip(): current_target.append(remainder[split_idx+1:].strip())
                        else:
                            if len(remainder) > 15: title = remainder[:12]; current_target.append(remainder[12:])
                            else: title = remainder
                continue
                
            if line_str.startswith('# '): title = line_str[2:]
            else: current_target.append(line_str)
                
        if not title and left_body: title = left_body[0][:15] 
        
        cleaned_left = [re.sub(r'^[\*\-]\s*', '', line) for line in left_body]
        cleaned_right = [re.sub(r'^[\*\-]\s*', '', line) for line in right_body]
            
        parsed_slides.append({
            'type': l_type, 'title': title, 'left_body': '\n'.join(cleaned_left), 'right_body': '\n'.join(cleaned_right), 'ai_prompt': ai_prompt
        })
    return parsed_slides

def get_english_translation(chinese_title):
    title_dict = {
        "背景": "PROJECT BACKGROUND", "需求": "REQUIREMENT ANALYSIS", "目标": "PROJECT OBJECTIVES", "参数": "TECHNICAL SPECIFICATIONS", 
        "工艺流程": "PROCESS FLOW", "工艺解析": "PROCESS ANALYSIS", "整体方案": "OVERALL SCHEME", "设备布局": "EQUIPMENT LAYOUT", 
        "核心机构": "CORE MECHANISM DESIGN", "结构设计": "STRUCTURAL DESIGN", "电气": "ELECTRICAL CONTROL SYSTEM", 
        "风险": "RISK ASSESSMENT (FMEA)", "计划": "PROJECT SCHEDULE", "清单": "BOM & SPARE PARTS", "结论": "DESIGN CONCLUSION"
    }
    for k, v in title_dict.items():
        if k in chinese_title: return v
    return "CHAPTER OVERVIEW"

def validate_api_key(api_key, provider, base_url=None):
    try:
        if "智谱" in provider:
            if "." not in api_key: return False, "❌ 格式错误：智谱 API Key 必须包含小数点。"
            client = ZhipuAI(api_key=api_key)
            client.chat.completions.create(model="glm-4-flash", messages=[{"role": "user", "content": "1"}], max_tokens=1)
            return True, "✅ 智谱引擎连接成功！"
        else:
            client = OpenAI(api_key=api_key, base_url=base_url)
            client.models.list()
            return True, "✅ 视觉引擎连接成功！"
    except Exception as e:
        return False, f"❌ 验证失败：{e}"

def create_ppt(parsed_slides, template_bytes, api_key, ai_provider, base_url=None, model_name=None):
    if not parsed_slides or not template_bytes: return None
    prs = Presentation(io.BytesIO(template_bytes))
    
    for i, slide_data in enumerate(parsed_slides):
        l_type = slide_data['type']
        try:
            if l_type == "封面":
                slide = prs.slides.add_slide(prs.slide_layouts[0])
                slide.placeholders[10].text = (slide_data['title'] + "\n" + slide_data['left_body']).strip()
            elif l_type == "目录":
                slide = prs.slides.add_slide(prs.slide_layouts[1])
                slide.placeholders[12].text = slide_data['left_body'].strip()
            elif l_type == "过渡":
                slide = prs.slides.add_slide(prs.slide_layouts[2])
                chinese_title = slide_data['title'].strip()
                slide.placeholders[10].text = chinese_title
                slide.placeholders[12].text = get_english_translation(chinese_title)
            elif l_type == "左文右图":
                slide = prs.slides.add_slide(prs.slide_layouts[3])
                slide.placeholders[10].text = slide_data['title']
                slide.placeholders[11].text = slide_data['left_body']
                
                if slide_data['ai_prompt'] and api_key:
                    st.toast(f"🎨 正在呼叫视觉引擎绘制概念图: {slide_data['title']} ...", icon="🤖")
                    img_url = ""
                    try:
                        if ai_provider == "官方 智谱 AI (CogView-3)":
                            client = ZhipuAI(api_key=api_key)
                            response = client.images.generations(model="cogview-3", prompt=slide_data['ai_prompt'])
                            img_url = response.data[0].url
                        elif ai_provider == "官方 OpenAI (DALL-E 3)":
                            client = OpenAI(api_key=api_key)
                            response = client.images.generate(model="dall-e-3", prompt=slide_data['ai_prompt'], n=1, size="1024x1024")
                            img_url = response.data[0].url
                        elif ai_provider == "通用聚合中转站 (支持 Midjourney/DALL-E等)":
                            client = OpenAI(api_key=api_key, base_url=base_url)
                            response = client.images.generate(model=model_name, prompt=slide_data['ai_prompt'], n=1, size="1024x1024")
                            img_url = response.data[0].url
                            
                        if img_url:
                            img_stream = io.BytesIO(requests.get(img_url).content)
                            slide.shapes.add_picture(img_stream, int(prs.slide_width * 0.55), int(prs.slide_height * 0.25), width=int(prs.slide_width * 0.40))
                    except Exception as draw_e:
                        pass
            
            elif l_type == "左右文字": 
                slide = prs.slides.add_slide(prs.slide_layouts[4])
                slide.placeholders[10].text = slide_data['title']
                slide.placeholders[11].text = slide_data['left_body']
                
                if 12 in [shape.placeholder_format.idx for shape in slide.placeholders]:
                    if slide_data['right_body']: slide.placeholders[12].text = slide_data['right_body']
                    else: slide.placeholders[12].text = "" 
                    
        except Exception as e:
            pass 
            
    try: prs.slides.add_slide(prs.slide_layouts[5])
    except: pass

    output = io.BytesIO()
    prs.save(output)
    output.seek(0)
    return output

# ================= UI 布局 =================
try: st.sidebar.image("JINYA-logo-透明背景.png", use_container_width=True)
except: st.sidebar.markdown("### 🏢 JINYA")

st.sidebar.markdown("<br>", unsafe_allow_html=True)
st.sidebar.markdown("#### 🛡️ 开发者授权配置")

with st.sidebar.form(key='api_form'):
    ai_provider = st.selectbox("🌐 视觉大模型引擎", ["官方 智谱 AI (CogView-3)", "官方 OpenAI (DALL-E 3)", "通用聚合中转站 (支持 Midjourney/DALL-E等)"])
    base_url, model_name = None, None
    if ai_provider == "通用聚合中转站 (支持 Midjourney/DALL-E等)":
        base_url = st.text_input("接口地址 (Base URL)", placeholder="https://api.xxx.com/v1")
        model_name = st.text_input("调用模型名称", value="dall-e-3")
    api_key_input = st.text_input("API 凭证 (API Key)", type="password")
    submit_btn = st.form_submit_button("✓ 保存配置", use_container_width=True)
    if submit_btn and api_key_input:
        with st.spinner("验证中..."):
            is_valid, msg = validate_api_key(api_key_input, ai_provider, base_url)
            if is_valid: st.success(msg)
            else: st.error(msg)

st.sidebar.markdown("""
<div style="cursor: pointer; color: #6B7280; font-size: 0.85rem; padding-top: 10px;">
    注：无需生成图片时，可不填写 API 凭证。
</div>
""", unsafe_allow_html=True)

st.markdown("<div class='main-title'>DFM生成引擎</div>", unsafe_allow_html=True)
st.markdown("<div class='sub-title'>自动化设备 DFM 极速演示方案生成器</div>", unsafe_allow_html=True)

# === 🌟 核心新增：新手使用指南 ===
with st.expander("💡 新手必看：只需 4 步，一键生成精美 DFM", expanded=False):
    st.markdown("""
    欢迎使用内部专属的 DFM 自动生成引擎！即便您是第一次使用，也能在 2 分钟内搞定专业排版。
    
    <span class="tutorial-step">📌 准备工作：让 AI 帮您写好文案。</span>
     点击下方蓝色的 **【📋 一键复制专属提示词】** 按钮。→ 打开您常用的大模型（如 豆包、Kimi、DeepSeek 或 ChatGPT），
     将复制的提示词发给它。→告诉大模型您当前项目的原始需求，它将为您输出一份**符合系统排版规则的标准 Markdown 文本**。
    
    <span class="tutorial-step"> ①：导入您的方案内容</span>
    您可以将大模型生成的方案保存为 Word（`.docx`）或 `.md` 文件上传，也可以直接复制纯文本，粘贴到 **方式二** 的输入框中。
    
    <span class="tutorial-step"> ②：上传官方 PPT 模板</span>
    将贵公司的标准汇报母版（`.pptx` 文件）上传至此处系统将自动识别母版排版格式。
    
    <span class="tutorial-step"> ③：配置自动画图（可选）</span>
    如果您希望系统在“左文右图”的页面中自动为您生成 **高清 3D 工业风概念图**，请在左侧栏填写您的视觉引擎 `API Key`，点击保存验证。如果只需纯文字排版，此步可跳过。
    
    <span class="tutorial-step"> ④：一键生成</span>
    点击最下方的 **【✨ 启动智能生成引擎 ✨】**。喝口水的功夫，一份图文并茂、双栏排版对齐、过渡页全英翻译的精美 PPT 就下载到您的电脑里啦！
    """, unsafe_allow_html=True)

# === 🎯 双栏排版专属提示词 ===
ai_prompt_text = """# Role (角色定位)
你是一位拥有10年经验的资深非标自动化机械工程师。你需要根据我提供的项目需求，撰写一份极其专业的《详细设计方案 (DFM)》。

# Output Format Strict Rules (输出格式绝对要求)
为了确保输出能够完美导入自动化渲染引擎，你必须【严格遵守】以下 Markdown 语法规则：

1. 分页符：每一页 PPT 之间，必须且只能用 `---` (三个减号) 独占一行来进行分隔。
2. 版式标签：每一页的开头第一行，必须指定版式，格式为 `[版式名称]`。
只能从以下 5 种版式中选择：`[封面]`, `[目录]`, `[过渡]`, `[左文右图]`, `[左右文字]`。
3. 页面标题：紧接版式标签的下一行，必须是标题，以 `# ` 开头。
4. 正文与双栏排版（核心指令！）：正文内容绝对禁止使用表格（|---|---|）。对于普通内容，直接使用圆点列表（*）排版。
如果你需要展示的内容较多，或者需要进行并列对比，请使用 `[左右文字]` 版式，
并且【必须】在该版式下使用 `[左]` 和 `[右]` 标签，将文字明确切分为左右两部分！
5. AI 配图指令：当版式为 [左文右图] 时，请在正文末尾单独起一行：[AI配图: 一段详细的中文画面描述]。
⚠️极度注意：你必须根据【当前页面的具体描述内容】
（如正在讲设备结构就写机械设备结构，讲工艺流程就写工艺流程），自己构思并生成相应的画面描述，生成图片要求工业风、3D概念图风格，并且要高清画质。
【绝对禁止】直接照抄下方示例模板里的配图提示词！

# 示例标准模板 (请大模型严格仿照此格式输出全文)：
[封面]
# 自动化设备详细设计方案 
---
[目录]
# 方案目录
* 一、项目背景与需求
* 二、整体方案及布局
---
[过渡]
# 二、整体方案及布局
---
[左文右图]
# 整体设备空间布局
* 设备尺寸预估为 2500x1800x2000 mm。
* 采用全封闭式钣金外壳。

[AI配图: 现代化自动化生产设备 3D 概念图，采用工业风、科技感图片，细节高清且精细]
---
[左右文字]
# 核心技术参数与优势
[左]
* 机械结构稳定，承载力高达500kg
* 运行节拍提升至 20s/pcs
* 占地面积节省 15%
[右]
* 全新一代 PLC 协同控制
* 支持一键式快速换型
* 兼容 MES/ERP 系统数据直连

# Task (任务)
请阅读上述规则。现在，我将提供新的项目信息，请按上述格式严格生成 Markdown 文件，并把这个MD文件以文件的格式发给我。不需要任何废话。"""

with st.expander("！！大模型发令枪：生成专业排版前必读！！", expanded=False):
    components.html(
        f"""<script>function copyToClipboard() {{ const text = {json.dumps(ai_prompt_text)}; navigator.clipboard.writeText(text).then(() => {{ const btn = document.getElementById("copy-btn"); btn.innerHTML = "✅ 复制成功！"; setTimeout(() => btn.innerHTML = "📋 一键复制专属提示词", 3000); }}); }}</script>
        <button id="copy-btn" onclick="copyToClipboard()" style="width: 100%; padding: 12px; background: #3B82F6; color: white; border: none; border-radius: 8px; font-weight: bold; cursor: pointer;">📋 一键复制专属提示词</button>""", height=60)
    st.code(ai_prompt_text, language="markdown")

st.markdown("#### Step 1. 导入您的方案内容")
with st.container(border=True):
    st.markdown("##### 📂 方式一：直接上传文档文件 (支持 DOCX / MD / TXT)")
    uploaded_doc = st.file_uploader("doc", type=["docx", "md", "txt"], label_visibility="collapsed")
    
    st.markdown("##### ✍️ 方式二：输入纯文本")
    raw_md_text = st.text_area("txt", height=250, label_visibility="collapsed", placeholder="如果没有文档，也可在此处直接粘贴/输入大模型生成的方案文本...")

st.markdown("#### Step 2. 上传 PPT 模板")
with st.container(border=True):
    uploaded_template = st.file_uploader("ppt", type=["pptx"], label_visibility="collapsed")

st.markdown("#### Step 3. 一键输出 DFM")
if st.button("✨ 启动智能生成引擎 ✨", use_container_width=True):
    
    md_text_to_process = ""
    if uploaded_doc:
        extracted = extract_text_content(uploaded_doc)
        if extracted: md_text_to_process = extracted
            
    if not md_text_to_process and raw_md_text:
        md_text_to_process = raw_md_text

    if uploaded_doc and not md_text_to_process and not raw_md_text:
        st.error("⚠️ 抱歉，当前系统暂不支持解析您上传的文件格式（如 PDF）。请上传 DOCX / MD 文件，或直接将文字粘贴在上方文本框中！")
    elif not md_text_to_process: 
        st.error("⚠️ 请导入内容！（上传有效文档或在文本框内粘贴文字）")
    elif not uploaded_template: 
        st.error("⚠️ 请在 Step 2 中上传 PPT 模板！")
    else:
        with st.spinner('🚀 引擎全速运转中，正在为您渲染精美排版...'):
            final_ppt = create_ppt(parse_markdown(md_text_to_process), uploaded_template.read(), api_key_input, ai_provider, base_url, model_name)
            if final_ppt:
                st.success("🎉 生成完毕！最纯净、最高效的双栏 DFM 报告已就绪。")
                st.download_button("📥 下载 PPT", final_ppt, "JINYA_DFM.pptx", use_container_width=True)